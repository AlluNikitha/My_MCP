import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Run 'pip install python-pptx' first.")
    sys.exit(1)

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Helper to apply pitch layout styles
    def apply_title_slide(slide, title_text, subtitle_text):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(8, 9, 12) # Dark cosmic black

        # Text frame
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title_text
        p1.font.name = 'Outfit'
        p1.font.size = Pt(54)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(99, 102, 241) # Indigo active tone
        p1.space_after = Pt(14)

        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Inter'
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(156, 163, 175) # Gray muted text

    def apply_content_slide(slide, title_text, bullets):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(17, 19, 25) # Charcoal secondary

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Outfit'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Bullets
        contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
        tf2 = contentBox.text_frame
        tf2.word_wrap = True

        for idx, item in enumerate(bullets):
            p_bullet = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
            p_bullet.text = item[0]
            p_bullet.font.name = 'Inter'
            p_bullet.font.size = Pt(20)
            p_bullet.font.bold = True
            p_bullet.font.color.rgb = RGBColor(99, 102, 241)
            p_bullet.space_after = Pt(4)
            p_bullet.space_before = Pt(8)

            p_desc = tf2.add_paragraph()
            p_desc.text = item[1]
            p_desc.font.name = 'Inter'
            p_desc.font.size = Pt(16)
            p_desc.font.color.rgb = RGBColor(209, 213, 219)
            p_desc.space_after = Pt(12)

    # 1. Slide 1: Title
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    apply_title_slide(slide1, "AURA: Smart Campus Operations Agent", "Chaining multi-domain Model Context Protocol (MCP) servers to automate college life.")

    # 2. Slide 2: Problem
    slide2 = prs.slides.add_slide(slide_layout)
    apply_content_slide(slide2, "The Problem: Fragmented Systems", [
        ("Dashboard Fatigue", "Existing college ERPs require users to navigate dozens of CRUD forms just to get simple answers or file operations."),
        ("Siloed Databases", "Attendance, Hostel, Library, and Finance are disconnected. Cross-domain policies require manual office coordination."),
        ("Chatbots Lacking Action", "Most AI chatbots are simple read-only question-answering systems that cannot execute updates or take operations.")
    ])

    # 3. Slide 3: Solution
    slide3 = prs.slides.add_slide(slide_layout)
    apply_content_slide(slide3, "The Solution: AURA Operations Agent", [
        ("Decoupled MCP Architecture", "Exposes all 9 campus operations domains as independent Model Context Protocol tool servers."),
        ("Multi-Domain Tool Chaining", "Agent orchestrator reasons across servers, calling multiple tools in sequence to fulfill queries."),
        ("Actionable Workflows", "Includes write-enabled operations (filing complaints, registering leave requests, applying for drives) rather than just reading tables.")
    ])

    # 4. Slide 4: Flagship Workflows
    slide4 = prs.slides.add_slide(slide_layout)
    apply_content_slide(slide4, "Flagship Demo Workflows", [
        ("Placement Eligibility Check (Flow A)", "Chains Attendance (check percentage) -> Finance (check fee dues) -> Library (check fines) -> Placements (verify CGPA criteria) in one step."),
        ("Weekend Hostel Outpass Approval (Flow B)", "Chains Hostel (get policy) -> Attendance (verify eligibility) -> Complaints (check disciplinary records) -> Hostel (issue approval + write outpass)."),
        ("Timetable watchdog conflict alert (Flow C)", "Proactive checks mapping student class slot timings against registered event slots, filing complaints automatically on clash detection.")
    ])

    # 5. Slide 5: Tech Stack
    slide5 = prs.slides.add_slide(slide_layout)
    apply_content_slide(slide5, "Architecture & Technology Stack", [
        ("React + Vite Frontend", "Features an interactive Agent Console trace log visualization and system-of-record views."),
        ("FastAPI Orchestrator", "Runs as the central client session layer, managing subprocess transport sessions to MCP servers and Anthropic LLM."),
        ("Model Context Protocol SDK", "FastMCP server implementations exposing schemas, JSON resources, and tools dynamically.")
    ])

    docs_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(docs_dir, "pitch_deck.pptx")
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at {output_path}")

if __name__ == "__main__":
    create_deck()
